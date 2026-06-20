from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_presentation():
    prs = Presentation()

    # Función auxiliar para títulos
    def add_title_slide(title_text, subtitle_text):
        slide_layout = prs.slide_layouts[0] 
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = title_text
        subtitle.text = subtitle_text
        return slide

    def add_bullet_slide(title_text, bullet_points):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        body_shape = slide.shapes.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(bullet_points):
            p = tf.add_paragraph()
            p.text = point
            p.level = 0
            p.font.size = Pt(24)
            # Agregar algo de espacio entre viñetas
            if i > 0:
                p.space_before = Pt(14)
                
        return slide

    # Diapositiva 1: Título
    add_title_slide(
        "DerivaShield", 
        "Aplicación del Cálculo Diferencial en Ciberseguridad Defensiva\nDetección de Anomalías de Red en Tiempo Real"
    )

    # Diapositiva 2: Introducción
    add_bullet_slide(
        "Objetivo del Proyecto",
        [
            "Diseñar un sistema de detección de intrusos (IDS) pasivo.",
            "Aplicar conceptos de Cálculo Diferencial en un entorno práctico.",
            "Detectar ataques de Denegación de Servicio (DDoS) y Escaneo de Puertos.",
            "Monitorear el tráfico de red en tiempo real mediante análisis matemático."
        ]
    )

    # Diapositiva 3: El Problema
    add_bullet_slide(
        "El Problema de la Detección Tradicional",
        [
            "Los sistemas tradicionales usan reglas fijas (ej. 'Alerta si hay más de 5000 paquetes').",
            "El tráfico de red moderno es altamente dinámico e impredecible.",
            "Reglas estáticas generan falsos positivos o fallan en detectar ataques lentos.",
            "Solución: Necesitamos medir la 'Tasa de Cambio' o aceleración del tráfico, no solo el volumen absoluto."
        ]
    )

    # Diapositiva 4: Aplicación Matemática
    add_bullet_slide(
        "La Solución: Cálculo Diferencial",
        [
            "Sea f(t) el volumen de tráfico (paquetes o puertos) en el instante t.",
            "La Primera Derivada f'(t) representa la tasa de cambio o aceleración instantánea.",
            "Como los datos son discretos (1 por segundo), usamos Diferencias Finitas:",
            "f'(t) ≈ [f(t) - f(t - Δt)] / Δt",
            "Un valor muy alto de f'(t) indica un pico anormal (posible ataque)."
        ]
    )

    # Diapositiva 5: Análisis Estadístico
    add_bullet_slide(
        "Umbrales Dinámicos Estadísticos",
        [
            "¿Cómo sabemos si f'(t) es peligrosamente alto?",
            "Calculamos la Media (μ) y la Desviación Estándar (σ) de las derivadas históricas.",
            "Establecemos un umbral dinámico: Umbral = μ + (k * σ).",
            "Donde 'k' es una constante de sensibilidad ajustable (Filtro Matemático).",
            "Si |f'(t)| > Umbral, el sistema dispara una alerta roja."
        ]
    )

    # Diapositiva 6: Casos de Uso
    add_bullet_slide(
        "Detección Específica",
        [
            "Ataques DDoS: Analizamos la métrica 'Paquetes por Segundo'. Un aumento abrupto (f'(t) alto) indica una inundación de datos.",
            "Port Scanning: Analizamos la métrica 'Puertos Destino Únicos'. Un escáner como Nmap contactará miles de puertos distintos rápidamente.",
            "Filtro de Ruido: Se requiere una aceleración mínima absoluta (ej. >100 paquetes/s²) para ignorar el ruido estadístico normal del sistema operativo."
        ]
    )

    # Diapositiva 7: Arquitectura de DerivaShield
    add_bullet_slide(
        "Arquitectura y Tecnologías",
        [
            "Lenguaje Core: Python 3.",
            "Captura de Red: Scapy y Npcap (Lectura de paquetes en Capa 2/3).",
            "Motor Matemático: Numpy (Vectores, Media, Desviación Estándar).",
            "Interfaz Gráfica (SOC): Dash y Plotly (Visualización en vivo, Dark Mode, Glassmorphism).",
            "Persistencia: SQLite3 (Registro histórico de alertas)."
        ]
    )

    # Diapositiva 8: Conclusión
    add_bullet_slide(
        "Conclusiones",
        [
            "El Cálculo Diferencial es una herramienta poderosa para analizar flujos de datos dinámicos.",
            "La tasa de cambio (Derivada) detecta ataques más rápido que los promedios estáticos.",
            "El sistema es puramente pasivo y no interfiere con la red, cumpliendo el rol defensivo.",
            "DerivaShield demuestra que las matemáticas puras tienen una aplicación directa en la Ciberseguridad moderna."
        ]
    )

    prs.save('DerivaShield_Presentacion.pptx')
    print("Presentación creada exitosamente.")

if __name__ == '__main__':
    create_presentation()
